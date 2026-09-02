#!/usr/bin/env python3
"""Export slides/deck_content.json to an editable 16:9 PPTX (native text, tables, embedded figures)."""
import json
import re
from html.parser import HTMLParser
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
content = json.loads((ROOT / "slides/deck_content.json").read_text(encoding="utf-8"))
BG, FG, ACC, SUB = RGBColor(0x0F, 0x14, 0x19), RGBColor(0xE6, 0xE9, 0xEF), RGBColor(0x5A, 0xA9, 0xFF), RGBColor(0x9F, 0xB3, 0xC8)


class P(HTMLParser):
    """Collect bullets (li at depth), tables (rows of cells), and pre blocks; drop tags."""
    def __init__(self):
        super().__init__(); self.bullets = []; self.tables = []; self.pre = []; self._li = None; self._depth = 0
        self._tab = None; self._row = None; self._cell = None; self._pre = None; self._big = []; self._in_big = 0

    def handle_starttag(self, t, a):
        a = dict(a)
        if t == "ul": self._depth += 1
        elif t == "li": self._li = []
        elif t == "table": self._tab = []
        elif t == "tr": self._row = []
        elif t in ("td", "th"): self._cell = []
        elif t == "pre": self._pre = []
        elif t == "div" and "big" in a.get("class", ""): self._in_big += 1; self._big.append([])
        elif t == "br" and self._in_big: self._big[-1].append("\n")
        elif t == "sub" or t == "sup": pass

    def handle_endtag(self, t):
        if t == "ul": self._depth -= 1
        elif t == "li" and self._li is not None:
            txt = re.sub(r"\s+", " ", "".join(self._li)).strip()
            if txt: self.bullets.append((max(0, self._depth - 1), txt))
            self._li = None
        elif t in ("td", "th") and self._cell is not None and self._row is not None:
            self._row.append(re.sub(r"\s+", " ", "".join(self._cell)).strip()); self._cell = None
        elif t == "tr" and self._row is not None and self._tab is not None:
            if self._row: self._tab.append(self._row)
            self._row = None
        elif t == "table" and self._tab is not None:
            self.tables.append(self._tab); self._tab = None
        elif t == "pre" and self._pre is not None:
            self.pre.append("".join(self._pre)); self._pre = None
        elif t == "div" and self._in_big: self._in_big -= 1

    def handle_data(self, d):
        if self._cell is not None: self._cell.append(d)
        elif self._li is not None: self._li.append(d)
        elif self._pre is not None: self._pre.append(d)
        elif self._in_big: self._big[-1].append(d)

    @property
    def big(self):
        return [re.sub(r"[ \t]+", " ", "".join(b)).strip() for b in self._big if "".join(b).strip()]


prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
blank = prs.slide_layouts[6]


def bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG


def text(slide, x, y, w, h, s, size=20, color=FG, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(s.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); r = p.add_run(); r.text = line
        r.font.size, r.font.color.rgb, r.font.bold = Pt(size), color, bold
    return tb


def bullets(slide, x, y, w, h, items, size=16):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame; tf.word_wrap = True
    for i, (lvl, t) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + t; r.font.size = Pt(size - 2 * lvl); r.font.color.rgb = FG
        p.space_after = Pt(6)


def table(slide, x, y, w, h, rows, size=11):
    nrow, ncol = len(rows), max(len(r) for r in rows)
    shp = slide.shapes.add_table(nrow, ncol, Inches(x), Inches(y), Inches(w), Inches(h)); tbl = shp.table
    for i, row in enumerate(rows):
        for j in range(ncol):
            c = tbl.cell(i, j); c.text = row[j] if j < len(row) else ""
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(size); r.font.bold = i == 0; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i == 0 else RGBColor(0x11, 0x11, 0x11)
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79) if i == 0 else (RGBColor(0xF4, 0xF6, 0xF8) if i % 2 else RGBColor(0xFF, 0xFF, 0xFF))


for sl in content["slides"]:
    s = prs.slides.add_slide(blank); bg(s)
    p = P(); p.feed(sl["body"])
    title = re.sub(r"<[^>]+>", " ", sl["title"]).replace("  ", " ").strip()
    if sl.get("kind") == "title":
        text(s, 0.8, 2.0, 11.5, 1.8, title, 40, RGBColor(0xFF, 0xFF, 0xFF), True)
        text(s, 0.8, 4.0, 11.5, 1.6, "\n".join(p.big), 20, FG)
        text(s, 0.8, 6.3, 11.5, 0.6, re.sub(r"<[^>]+>", "", sl.get("sub", "")), 12, SUB)
        continue
    text(s, 0.6, 0.35, 12.1, 0.8, title, 28, RGBColor(0xFF, 0xFF, 0xFF), True)
    if sl.get("sub"):
        text(s, 0.6, 1.05, 12.1, 0.4, re.sub(r"<[^>]+>", "", sl["sub"]), 12, SUB)
    figs = list(sl.get("figs", {}).values())
    has_fig = bool(figs)
    left_w = 6.4 if has_fig else 12.1
    y = 1.55
    if p.big and not p.bullets and not p.tables:
        text(s, 0.6, y, left_w, 5.3, "\n\n".join(p.big), 18)
    if p.tables:
        for t in p.tables[:1]:
            table(s, 0.6, y, left_w, min(5.2, 0.32 * len(t) + 0.3), t, 11 if len(t) < 9 else 9)
            y += min(5.2, 0.32 * len(t) + 0.3) + 0.2
    if p.bullets:
        bullets(s, 0.6, y, left_w, 7.0 - y, p.bullets, 16 if len(p.bullets) <= 6 else 14)
    if p.pre:
        text(s, 0.6 if not p.bullets else 7.2, 1.55, 5.6, 5.0, p.pre[0], 10, FG)
    if has_fig:
        s.shapes.add_picture(str(ROOT / figs[0]), Inches(7.2), Inches(1.55), width=Inches(5.7))
    text(s, 0.6, 7.05, 12.1, 0.3, content["footer"], 10, SUB)
out = ROOT / "slides/awesome_diffusion_OT_deck.pptx"
prs.save(out); print("wrote", out, f"{out.stat().st_size//1024} KB, {len(prs.slides)} slides")
